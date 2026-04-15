import json
from pathlib import Path
import pdfplumber
from docx import Document
from pptx import Presentation
import email
from email import policy

#Config
CHUNK_SIZE = 200
OVERLAP = 30

#Chunking
def chunk_text(text, size=CHUNK_SIZE, overlap=OVERLAP):
    words = text.split()
    chunks = []

    for i in range(0, len(words), size - overlap):
        chunk = " ".join(words[i:i + size])
        if len(chunk.strip()) > 30:
            chunks.append(chunk)

    return chunks

#File Parsers
#PDF
def parse_pdf(file_path):
    text = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text.append(t)
    return "\n".join(text)

#Docx
def parse_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs if p.text])

#PPTX
def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)

#Email
def parse_eml(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        msg = email.message_from_file(f, policy=policy.default)

    if msg.is_multipart():
        parts = []
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                parts.append(part.get_content())
        return "\n".join(parts)
    else:
        return msg.get_content()

#JSON Glossary
def parse_json_glossary(file_path):
    glossary = {}

    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    for key, val in data.get("V", {}).items():
        definition = val.get("definition", "")
        if definition:
            glossary[key] = definition

    return glossary

#Loader
def load_documents(base_path="test_files"):
    documents = []
    glossary = {}

    base = Path(base_path)

    for file in base.iterdir():
        suffix = file.suffix.lower()

        try:
            if suffix == ".json":
                glossary.update(parse_json_glossary(file))

            elif suffix == ".pdf":
                text = parse_pdf(file)
                for chunk in chunk_text(text):
                    documents.append({"text": chunk, "source": file.name})

            elif suffix == ".docx":
                text = parse_docx(file)
                for chunk in chunk_text(text):
                    documents.append({"text": chunk, "source": file.name})

            elif suffix == ".pptx":
                text = parse_pptx(file)
                for chunk in chunk_text(text):
                    documents.append({"text": chunk, "source": file.name})

            elif suffix == ".eml":
                text = parse_eml(file)
                for chunk in chunk_text(text):
                    documents.append({"text": chunk, "source": file.name})

        except Exception as e:
            print(f"[WARN] Failed to parse {file.name}: {e}")

    return documents, glossary